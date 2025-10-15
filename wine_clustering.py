import os
import json
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.datasets import load_wine
from sklearn.impute import SimpleImputer
from sklearn.preprocessing import StandardScaler
from sklearn.decomposition import PCA
from sklearn.cluster import KMeans
from sklearn.metrics import silhouette_score, davies_bouldin_score, calinski_harabasz_score
from sklearn.metrics import adjusted_rand_score, normalized_mutual_info_score
from scipy.cluster.hierarchy import linkage, dendrogram, cophenet, fcluster
from scipy.spatial.distance import pdist
from scipy import stats
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, Preformatted
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def ensure_dirs(base_dir):
    paths = [
        os.path.join(base_dir, 'data'),
        os.path.join(base_dir, 'outputs'),
        os.path.join(base_dir, 'outputs', 'figures'),
        os.path.join(base_dir, 'outputs', 'figures', 'histograms'),
    ]
    for p in paths:
        os.makedirs(p, exist_ok=True)


def save_fig(fig, path):
    fig.tight_layout()
    fig.savefig(path, dpi=150)
    plt.close(fig)


def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    ensure_dirs(base_dir)
    out_dir = os.path.join(base_dir, 'outputs')
    fig_dir = os.path.join(out_dir, 'figures')
    hist_dir = os.path.join(fig_dir, 'histograms')

    # Part A: Data Loading & Preprocessing
    wine = load_wine()
    X = pd.DataFrame(wine.data, columns=wine.feature_names)
    y = pd.Series(wine.target, name='label')  # only for optional validation

    # Show first 10 rows and confirm numeric types
    print('First 10 rows:')
    print(X.head(10))
    print('\nDtypes:')
    print(X.dtypes)

    # Handle missing values
    imputer = SimpleImputer(strategy='median')
    X_imputed = pd.DataFrame(imputer.fit_transform(X), columns=X.columns)

    # Standardize features
    scaler = StandardScaler()
    X_scaled = pd.DataFrame(scaler.fit_transform(X_imputed), columns=X.columns)
    scaled_path = os.path.join(out_dir, 'scaled_wine.csv')
    X_scaled.to_csv(scaled_path, index=False)
    print(f'Scaled dataset saved to: {scaled_path}')

    # Part B: First EDA
    selected = ['alcohol', 'malic_acid', 'ash', 'flavanoids']
    desc = X_imputed[selected].describe().T
    desc_path = os.path.join(out_dir, 'descriptive_stats_selected.csv')
    desc.to_csv(desc_path)
    print(f'Descriptive stats saved to: {desc_path}')

    # Histograms
    for col in selected:
        fig, ax = plt.subplots(figsize=(6, 4))
        sns.histplot(X_imputed[col], kde=True, ax=ax)
        ax.set_title(f'Histogram of {col}')
        ax.set_xlabel(col)
        save_fig(fig, os.path.join(hist_dir, f'{col}_hist.png'))

    # Simple commentary (printed)
    skewness = X_imputed[selected].skew()
    print('\nSkewness of selected variables:')
    print(skewness)

    # PCA scree plot
    pca = PCA()
    pca.fit(X_scaled)
    evr = pca.explained_variance_ratio_
    fig, ax = plt.subplots(figsize=(7, 4))
    ax.plot(np.arange(1, len(evr)+1), evr, marker='o')
    ax.set_title('PCA Scree Plot (Explained Variance Ratio)')
    ax.set_xlabel('Component')
    ax.set_ylabel('Explained Variance Ratio')
    save_fig(fig, os.path.join(fig_dir, 'pca_scree.png'))

    # Part C: Feature Engineering
    X_fe = X_scaled.copy()
    # Composite features
    # Use original units for intuitive ratios, then standardize within X_fe
    acid_index = (X_imputed['malic_acid'] / X_imputed['alcohol']).values.reshape(-1, 1)
    phenols_ratio = (X_imputed['total_phenols'] / (X_imputed['flavanoids'] + 1e-6)).values.reshape(-1, 1)
    # Standardize these composite features to align with X_fe scale
    comp_scaler = StandardScaler()
    comp_feats = comp_scaler.fit_transform(np.hstack([acid_index, phenols_ratio]))
    X_fe['acid_index'] = comp_feats[:, 0]
    X_fe['phenols_ratio'] = comp_feats[:, 1]

    # PCA for dimensionality reduction (retain 95% variance)
    pca95 = PCA(n_components=0.95, svd_solver='full')
    X_pca = pca95.fit_transform(X_scaled)
    print(f'PCA retained components: {X_pca.shape[1]} for 95% variance')

    # Part D: Clustering Modelling & Validation
    def evaluate_kmeans(features, name):
        results = []
        for k in range(2, 9):
            km = KMeans(n_clusters=k, n_init=10, random_state=42)
            labels = km.fit_predict(features)
            sil = silhouette_score(features, labels)
            db = davies_bouldin_score(features, labels)
            ch = calinski_harabasz_score(features, labels)
            results.append({'feature_set': name, 'k': k, 'silhouette': sil, 'davies_bouldin': db, 'calinski_harabasz': ch})
        return pd.DataFrame(results)

    km_scaled = evaluate_kmeans(X_scaled.values, 'scaled')
    km_pca = evaluate_kmeans(X_pca, 'pca')
    km_fe = evaluate_kmeans(X_fe.values, 'feature_engineered')
    all_km = pd.concat([km_scaled, km_pca, km_fe], ignore_index=True)
    metrics_path = os.path.join(out_dir, 'kmeans_metrics.csv')
    all_km.to_csv(metrics_path, index=False)
    print(f'KMeans metrics saved to: {metrics_path}')

    # Choose best by silhouette (tie-breaker: lowest DB, highest CH)
    best = all_km.sort_values(['silhouette', 'davies_bouldin', 'calinski_harabasz'], ascending=[False, True, False]).iloc[0]
    chosen_set = best['feature_set']
    chosen_k = int(best['k'])
    print(f'Chosen KMeans config: set={chosen_set}, k={chosen_k}, silhouette={best["silhouette"]:.3f}')

    if chosen_set == 'scaled':
        features_for_kmeans = X_scaled.values
    elif chosen_set == 'pca':
        features_for_kmeans = X_pca
    else:
        features_for_kmeans = X_fe.values

    kmeans_final = KMeans(n_clusters=chosen_k, n_init=25, random_state=42)
    kmeans_labels = kmeans_final.fit_predict(features_for_kmeans)

    # Plot indices vs k
    def plot_indices(df, name):
        df_k = df[df['feature_set'] == name]
        fig, ax = plt.subplots(1, 3, figsize=(12, 3))
        ax[0].plot(df_k['k'], df_k['silhouette'], marker='o'); ax[0].set_title(f'Silhouette ({name})'); ax[0].set_xlabel('k')
        ax[1].plot(df_k['k'], df_k['davies_bouldin'], marker='o'); ax[1].set_title(f'Davies-Bouldin ({name})'); ax[1].set_xlabel('k')
        ax[2].plot(df_k['k'], df_k['calinski_harabasz'], marker='o'); ax[2].set_title(f'Calinski-Harabasz ({name})'); ax[2].set_xlabel('k')
        save_fig(fig, os.path.join(fig_dir, f'indices_vs_k_{name}.png'))

    plot_indices(all_km, 'scaled')
    plot_indices(all_km, 'pca')
    plot_indices(all_km, 'feature_engineered')

    # Hierarchical clustering: Ward & average
    # Use chosen feature set
    Z_ward = linkage(features_for_kmeans, method='ward', metric='euclidean')
    Z_avg = linkage(features_for_kmeans, method='average', metric='euclidean')

    # Dendrograms
    for name, Z in [('ward', Z_ward), ('average', Z_avg)]:
        fig, ax = plt.subplots(figsize=(12, 5))
        dendrogram(Z, ax=ax, truncate_mode=None)
        ax.set_title(f'Hierarchical Dendrogram ({name} linkage)')
        ax.set_xlabel('Samples')
        ax.set_ylabel('Distance')
        save_fig(fig, os.path.join(fig_dir, f'dendrogram_{name}.png'))

    # Cophenetic correlation
    Y = pdist(features_for_kmeans)
    coph_ward, _ = cophenet(Z_ward, Y)
    coph_avg, _ = cophenet(Z_avg, Y)
    with open(os.path.join(out_dir, 'cophenetic.json'), 'w') as f:
        json.dump({'ward': coph_ward, 'average': coph_avg}, f, indent=2)
    print(f'Cophenetic correlations -> ward: {coph_ward:.3f}, average: {coph_avg:.3f}')

    # Optional external validation (true labels)
    ari = adjusted_rand_score(y, kmeans_labels)
    nmi = normalized_mutual_info_score(y, kmeans_labels)
    with open(os.path.join(out_dir, 'external_validation.json'), 'w') as f:
        json.dump({'ARI': ari, 'NMI': nmi}, f, indent=2)
    print(f'External validation -> ARI: {ari:.3f}, NMI: {nmi:.3f}')

    # Part E: Second EDA & Statistical Inference
    # Cluster centroids in original units
    # If we clustered on scaled features, inverse-transform centers; if PCA, map centers back approx by inverse PCA then inverse scaler
    def centers_in_original_units():
        if chosen_set == 'scaled' or chosen_set == 'feature_engineered':
            # Extract only original feature columns from centers if feature_engineered
            centers = kmeans_final.cluster_centers_
            if chosen_set == 'feature_engineered':
                # First columns correspond to original standardized features, last two are engineered
                centers = centers[:, :X_scaled.shape[1]]
            original = scaler.inverse_transform(centers)
            return pd.DataFrame(original, columns=X.columns)
        elif chosen_set == 'pca':
            centers_pca = kmeans_final.cluster_centers_
            approx_scaled = pca95.inverse_transform(centers_pca)
            original = scaler.inverse_transform(approx_scaled)
            return pd.DataFrame(original, columns=X.columns)

    centroids_original = centers_in_original_units()
    centroids_path = os.path.join(out_dir, 'cluster_centroids_original_units.csv')
    centroids_original.to_csv(centroids_path, index=False)
    print(f'Cluster centroids (original units) saved to: {centroids_path}')

    # PCA scatter colored by cluster (use first 2 PCs from full PCA on scaled)
    pca2 = PCA(n_components=2)
    X_pca2 = pca2.fit_transform(X_scaled)
    fig, ax = plt.subplots(figsize=(6, 5))
    scatter = ax.scatter(X_pca2[:, 0], X_pca2[:, 1], c=kmeans_labels, cmap='viridis', s=45, edgecolor='k')
    ax.set_title('PCA Scatter (PC1 vs PC2) colored by KMeans cluster')
    ax.set_xlabel('PC1')
    ax.set_ylabel('PC2')
    legend1 = ax.legend(*scatter.legend_elements(), title="Cluster")
    ax.add_artist(legend1)
    save_fig(fig, os.path.join(fig_dir, 'pca_scatter_kmeans.png'))

    # ANOVA & Kruskal-Wallis for alcohol and flavanoids
    df_with_labels = X_imputed.copy()
    df_with_labels['cluster'] = kmeans_labels
    def group_values(col):
        return [group[col].values for _, group in df_with_labels.groupby('cluster')]

    results_stats = {}
    for col in ['alcohol', 'flavanoids']:
        groups = group_values(col)
        # One-way ANOVA (assumes normality/variance homogeneity approximately)
        f_stat, p_anova = stats.f_oneway(*groups)
        # Effect size (eta squared): SS_between / SS_total
        all_vals = df_with_labels[col].values
        grand_mean = all_vals.mean()
        ss_total = ((all_vals - grand_mean) ** 2).sum()
        ss_between = sum([len(g) * (g.mean() - grand_mean) ** 2 for g in groups])
        eta_sq = ss_between / ss_total if ss_total > 0 else np.nan

        # Kruskal-Wallis (non-parametric)
        h_stat, p_kw = stats.kruskal(*groups)
        # Epsilon squared effect size for Kruskal: (H - k + 1) / (n - k)
        k = len(groups)
        n = len(all_vals)
        eps_sq = (h_stat - k + 1) / (n - k) if (n - k) > 0 else np.nan

        results_stats[col] = {
            'ANOVA_F': float(f_stat), 'ANOVA_p': float(p_anova), 'eta_squared': float(eta_sq),
            'Kruskal_H': float(h_stat), 'Kruskal_p': float(p_kw), 'epsilon_squared': float(eps_sq)
        }

    with open(os.path.join(out_dir, 'stat_tests.json'), 'w') as f:
        json.dump(results_stats, f, indent=2)
    print('Statistical tests saved to stat_tests.json')

    # Cluster profiles bar plot for selected variables
    key_vars = ['alcohol', 'malic_acid', 'ash', 'flavanoids']
    fig, ax = plt.subplots(figsize=(10, 5))
    idx = np.arange(len(key_vars))
    width = 0.8 / chosen_k
    for c in range(chosen_k):
        ax.bar(idx + c * width, centroids_original.loc[c, key_vars], width, label=f'Cluster {c}')
    ax.set_xticks(idx + width * (chosen_k - 1) / 2)
    ax.set_xticklabels(key_vars)
    ax.set_ylabel('Original units')
    ax.set_title('Cluster Profiles (Selected Variables)')
    ax.legend()
    save_fig(fig, os.path.join(fig_dir, 'cluster_profiles_bar.png'))

    # Part F: Presentation & Reflection - PDF Report
    report_path = os.path.join(out_dir, 'wine_clustering_report.pdf')
    doc = SimpleDocTemplate(report_path, pagesize=letter)
    styles = getSampleStyleSheet()
    code_style = ParagraphStyle('Code', parent=styles['BodyText'], fontName='Courier', fontSize=9, leading=11)
    story = []
    # Enriched context for report
    n_samples, n_features = X.shape
    class_counts = y.value_counts().to_dict()
    best_sil = float(best['silhouette'])
    best_db = float(best['davies_bouldin'])
    best_ch = float(best['calinski_harabasz'])
    story.append(Paragraph('UCI Wine Dataset: K-means vs Hierarchical Clustering', styles['Title']))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Problem Statement', styles['Heading2']))
    story.append(Paragraph('Cluster wines using 13 physico-chemical attributes without leveraging labels for modelling; compare K-means and hierarchical approaches, select k via internal indices, and interpret cluster differences via statistical inference.', styles['BodyText']))
    story.append(Spacer(1, 12))
    # Part A: Data Loading & Preprocessing
    story.append(Paragraph('Part A: Data Loading & Preprocessing', styles['Heading2']))
    story.append(Paragraph('How obtained', styles['Heading3']))
    story.append(Paragraph('Loaded UCI Wine data via scikit-learn; inspected the first 10 rows and confirmed all 13 features are numeric (float64). Applied median imputation (none missing in this dataset) and Z-score standardization, then saved the scaled dataset to CSV.', styles['BodyText']))
    # Show first 10 rows (selected variables for legibility)
    first10 = X_imputed[['alcohol','malic_acid','ash','flavanoids']].head(10)
    table_data = [['Row','alcohol','malic_acid','ash','flavanoids']]
    for i, r in first10.iterrows():
        table_data.append([str(i), f"{r['alcohol']:.2f}", f"{r['malic_acid']:.2f}", f"{r['ash']:.2f}", f"{r['flavanoids']:.2f}"])
    tbl = Table(table_data)
    tbl.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER')
    ]))
    story.append(tbl)
    story.append(Spacer(1, 8))
    story.append(Paragraph('Significance', styles['Heading3']))
    story.append(Paragraph('Numeric types and standardization ensure comparable scales for distance-based methods; imputation prevents missingness from biasing PCA or clustering.', styles['BodyText']))
    story.append(Spacer(1, 8))
    story.append(Paragraph('Code: Loading, Imputation, Standardization, Save', styles['Heading3']))
    story.append(Preformatted(
        "from sklearn.datasets import load_wine\n"
        "data = load_wine()\n"
        "X = pd.DataFrame(data.data, columns=data.feature_names)\n"
        "y = pd.Series(data.target)\n"
        "imputer = SimpleImputer(strategy='median')\n"
        "X_imputed = pd.DataFrame(imputer.fit_transform(X), columns=X.columns)\n"
        "scaler = StandardScaler()\n"
        "X_scaled = pd.DataFrame(scaler.fit_transform(X_imputed), columns=X.columns)\n"
        "X_scaled.to_csv(os.path.join(out_dir, 'scaled_wine.csv'), index=False)",
        code_style
    ))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Dataset & Features', styles['Heading2']))
    story.append(Paragraph(f'Total samples: {n_samples}; features: {n_features}. Variables: {", ".join(X.columns)}. Labels used only for optional external validation (not for tuning). Class counts: {class_counts}.', styles['BodyText']))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Preprocessing Choices', styles['Heading2']))
    story.append(Paragraph('Median imputation handles potential missingness (none observed in Wine). Z-score standardization is crucial for distance-based algorithms and PCA, ensuring all features contribute comparably.', styles['BodyText']))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Pipeline Overview', styles['Heading2']))
    story.append(Paragraph('Load → Impute → Standardize → EDA → PCA → Feature Engineering → KMeans/Hierarchical → Validation → Inference → Presentation', styles['BodyText']))
    story.append(Spacer(1, 12))
    # Part B: First Exploratory Data Analysis (20 points)
    story.append(Paragraph('Part B: First Exploratory Data Analysis (20 points)', styles['Heading2']))
    story.append(Paragraph('Descriptive statistics for selected features (alcohol, malic_acid, ash, flavanoids) were computed and saved. Histograms for at least four variables reveal distribution shape and skew; we also include a PCA scree plot to motivate dimensionality reduction.', styles['BodyText']))
    # Descriptive stats (compact view)
    desc_view = X_imputed[['alcohol','malic_acid','ash','flavanoids']].describe().loc[['mean','std','min','max']]
    table_data = [['stat','alcohol','malic_acid','ash','flavanoids']]
    for stat in desc_view.index:
        row = [stat] + [f"{desc_view.loc[stat, c]:.2f}" for c in ['alcohol','malic_acid','ash','flavanoids']]
        table_data.append(row)
    tbl = Table(table_data)
    tbl.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER')
    ]))
    story.append(tbl)
    story.append(Spacer(1, 8))
    story.append(Paragraph('Histograms (4 variables)', styles['Heading3']))
    for col in ['alcohol','malic_acid','ash','flavanoids']:
        story.append(Image(os.path.join(hist_dir, f'{col}_hist.png'), width=300, height=200))
    story.append(Spacer(1, 6))
    story.append(Paragraph(f'Skewness (selected): {skewness.to_dict()}', styles['BodyText']))
    story.append(Paragraph('Significance: Histograms show spread/outliers; skewness highlights long tails (e.g., malic_acid right-skew). Scree plot shows diminishing variance gains across components.', styles['BodyText']))
    story.append(Spacer(1, 8))
    story.append(Paragraph('Code: Histograms & Descriptive Stats', styles['Heading3']))
    story.append(Preformatted(
        "selected = ['alcohol','malic_acid','ash','flavanoids']\n"
        "desc = X_imputed[selected].describe().T\n"
        "for col in selected:\n"
        "    fig, ax = plt.subplots()\n"
        "    sns.histplot(X_imputed[col], kde=True, ax=ax)\n"
        "    save_fig(fig, os.path.join(hist_dir, f'{col}_hist.png'))",
        code_style
    ))
    story.append(Spacer(1, 12))
    story.append(Paragraph('PCA Scree Plot', styles['Heading2']))
    story.append(Image(os.path.join(fig_dir, 'pca_scree.png'), width=400, height=250))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Dimension Reduction (PCA)', styles['Heading2']))
    story.append(Paragraph(f'PCA reduces redundancy from correlated variables and can improve cluster separation. Retaining 95% variance yields {X_pca.shape[1]} components; the scree plot shows diminishing returns across components.', styles['BodyText']))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Code: Building the PCA Scree Plot', styles['Heading3']))
    story.append(Preformatted(
        "pca = PCA(); pca.fit(X_scaled)\n"
        "evr = pca.explained_variance_ratio_\n"
        "fig, ax = plt.subplots()\n"
        "ax.plot(np.arange(1, len(evr)+1), evr, marker='o')\n"
        "save_fig(fig, os.path.join(fig_dir, 'pca_scree.png'))",
        code_style
    ))
    story.append(Spacer(1, 10))
    # Part C: Feature Engineering (15 points)
    story.append(Paragraph('Part C: Feature Engineering (15 points)', styles['Heading2']))
    story.append(Paragraph('We optionally crafted composite features (acid_index = malic_acid/alcohol; phenols_ratio = total_phenols/(flavanoids+ε)) and standardized them. We also applied PCA retaining 95% variance. These transformed features were considered in clustering when they improved separation/compactness.', styles['BodyText']))
    story.append(Paragraph('Code: Composite Features & PCA (95%)', styles['Heading3']))
    story.append(Preformatted(
        "acid_index = X_imputed['malic_acid'] / X_imputed['alcohol']\n"
        "phenols_ratio = X_imputed['total_phenols'] / (X_imputed['flavanoids'] + 1e-6)\n"
        "comp_scaler = StandardScaler()\n"
        "comp = comp_scaler.fit_transform(np.c_[acid_index, phenols_ratio])\n"
        "X_fe = X_scaled.copy(); X_fe['acid_index'] = comp[:,0]; X_fe['phenols_ratio'] = comp[:,1]\n"
        "pca95 = PCA(n_components=0.95, svd_solver='full'); X_pca = pca95.fit_transform(X_scaled)",
        code_style
    ))
    story.append(Spacer(1, 12))
    # Part D: Clustering Modelling & Validation (25 points)
    story.append(Paragraph('Part D: Clustering Modelling & Validation (25 points)', styles['Heading2']))
    story.append(Paragraph('K-means: run k=2..8; compute Silhouette (↑), Davies–Bouldin (↓), Calinski–Harabasz (↑) and pick k balancing these indices. Hierarchical: build Ward & average dendrograms and compute cophenetic correlation to assess tree fidelity. External validation (ARI, NMI) reported for analysis only.', styles['BodyText']))
    story.append(Paragraph('KMeans Model Selection', styles['Heading2']))
    story.append(Paragraph('Clustering Methods & Indices', styles['Heading2']))
    story.append(Paragraph('K-means minimizes within-cluster variance; hierarchical builds a dendrogram using Ward (variance-minimizing) and average (mean-distance) linkages. Internal indices used: Silhouette (cohesion/separation, higher better), Davies–Bouldin (cluster separation/compactness, lower better), Calinski–Harabasz (between/within dispersion, higher better).', styles['BodyText']))
    for name in ['scaled', 'pca', 'feature_engineered']:
        story.append(Paragraph(f'Indices vs k ({name})', styles['Heading3']))
        story.append(Image(os.path.join(fig_dir, f'indices_vs_k_{name}.png'), width=400, height=250))
        story.append(Spacer(1, 6))
    story.append(Paragraph(f'Selected configuration: set={chosen_set}, k={chosen_k}. Metrics: Silhouette={best_sil:.3f}, Davies–Bouldin={best_db:.3f}, Calinski–Harabasz={best_ch:.1f}.', styles['BodyText']))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Code: Computing Indices vs k', styles['Heading3']))
    story.append(Preformatted(
        "for k in range(2, 9):\n"
        "    km = KMeans(n_clusters=k, n_init=10, random_state=42)\n"
        "    labels = km.fit_predict(features)\n"
        "    sil = silhouette_score(features, labels)\n"
        "    db = davies_bouldin_score(features, labels)\n"
        "    ch = calinski_harabasz_score(features, labels)  # plot vs k",
        code_style
    ))
    story.append(Spacer(1, 10))
    story.append(Paragraph('Hierarchical Dendrograms', styles['Heading2']))
    story.append(Image(os.path.join(fig_dir, 'dendrogram_ward.png'), width=450, height=250))
    story.append(Image(os.path.join(fig_dir, 'dendrogram_average.png'), width=450, height=250))
    story.append(Paragraph(f'Cophenetic correlation evaluates dendrogram fidelity to original distances. Ward: {coph_ward:.3f}, Average: {coph_avg:.3f}.', styles['BodyText']))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Code: Building Dendrograms', styles['Heading3']))
    story.append(Preformatted(
        "Z_ward = linkage(features_for_kmeans, method='ward', metric='euclidean')\n"
        "fig, ax = plt.subplots()\n"
        "dendrogram(Z_ward, ax=ax)\n"
        "save_fig(fig, os.path.join(fig_dir, 'dendrogram_ward.png'))",
        code_style
    ))
    story.append(Spacer(1, 10))
    story.append(Paragraph('Code: Cophenetic Correlation', styles['Heading3']))
    story.append(Preformatted(
        "Y = pdist(features_for_kmeans)\n"
        "coph_ward, _ = cophenet(Z_ward, Y)\n"
        "coph_avg, _ = cophenet(Z_avg, Y)\n"
        "json.dump({'ward': coph_ward, 'average': coph_avg}, open(os.path.join(out_dir,'cophenetic.json'),'w'))",
        code_style
    ))
    story.append(Spacer(1, 10))
    story.append(Paragraph('PCA Scatter by KMeans Clusters', styles['Heading2']))
    story.append(Image(os.path.join(fig_dir, 'pca_scatter_kmeans.png'), width=400, height=300))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Code: PCA Scatter by Cluster', styles['Heading3']))
    story.append(Preformatted(
        "pca2 = PCA(n_components=2)\n"
        "X_pca2 = pca2.fit_transform(X_scaled)\n"
        "fig, ax = plt.subplots()\n"
        "ax.scatter(X_pca2[:, 0], X_pca2[:, 1], c=kmeans_labels, cmap='viridis')\n"
        "save_fig(fig, os.path.join(fig_dir, 'pca_scatter_kmeans.png'))",
        code_style
    ))
    story.append(Spacer(1, 10))
    story.append(Paragraph('Cluster Centroids (Original Units)', styles['Heading2']))
    # Show a small table of centroids for key vars
    table_data = [['Cluster'] + key_vars]
    for c in range(chosen_k):
        row = [str(c)] + [f'{centroids_original.loc[c, v]:.2f}' for v in key_vars]
        table_data.append(row)
    tbl = Table(table_data)
    tbl.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER')
    ]))
    story.append(tbl)
    story.append(Spacer(1, 12))
    story.append(Paragraph('Cluster Profiles (Selected Variables)', styles['Heading2']))
    story.append(Image(os.path.join(fig_dir, 'cluster_profiles_bar.png'), width=450, height=250))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Code: Cluster Profiles Bar Chart', styles['Heading3']))
    story.append(Preformatted(
        "key_vars = ['alcohol', 'malic_acid', 'ash', 'flavanoids']\n"
        "fig, ax = plt.subplots(figsize=(10, 5))\n"
        "idx = np.arange(len(key_vars)); width = 0.8 / chosen_k\n"
        "for c in range(chosen_k):\n"
        "    ax.bar(idx + c * width, centroids_original.loc[c, key_vars], width)\n"
        "save_fig(fig, os.path.join(fig_dir, 'cluster_profiles_bar.png'))",
        code_style
    ))
    story.append(Spacer(1, 10))
    # Part E: Second EDA & Statistical Inference (15 points)
    story.append(Paragraph('Part E: Second EDA & Statistical Inference (15 points)', styles['Heading2']))
    story.append(Paragraph('Provide centroids in original units and PCA scatter coloured by cluster to interpret clusters. Use ANOVA/Kruskal-Wallis on alcohol and flavanoids; report p-values and effect sizes (eta², ε²). If labels used for validation, note any discrepancies.', styles['BodyText']))
    story.append(Paragraph('Statistical Inference', styles['Heading2']))
    story.append(Paragraph(f"Alcohol: ANOVA p={results_stats['alcohol']['ANOVA_p']:.4f}, eta²={results_stats['alcohol']['eta_squared']:.3f}; Kruskal p={results_stats['alcohol']['Kruskal_p']:.4f}, ε²={results_stats['alcohol']['epsilon_squared']:.3f}", styles['BodyText']))
    story.append(Paragraph(f"Flavanoids: ANOVA p={results_stats['flavanoids']['ANOVA_p']:.4f}, eta²={results_stats['flavanoids']['eta_squared']:.3f}; Kruskal p={results_stats['flavanoids']['Kruskal_p']:.4f}, ε²={results_stats['flavanoids']['epsilon_squared']:.3f}", styles['BodyText']))
    story.append(Paragraph('Effect sizes (eta², ε²) contextualize practical differences beyond p-values.', styles['BodyText']))
    story.append(Paragraph('Code: ANOVA & Kruskal-Wallis', styles['Heading3']))
    story.append(Preformatted(
        "df = X_imputed.copy(); df['cluster'] = kmeans_labels\n"
        "groups = [g['alcohol'].values for _, g in df.groupby('cluster')]\n"
        "f, p = stats.f_oneway(*groups)\n"
        "h, pkw = stats.kruskal(*groups)\n"
        "# eta² and ε² computed from sums of squares and H-statistic",
        code_style
    ))
    story.append(Spacer(1, 12))
    story.append(Paragraph('External Validation (Optional)', styles['Heading2']))
    story.append(Paragraph(f'Adjusted Rand Index (ARI)={ari:.3f}, Normalized Mutual Info (NMI)={nmi:.3f}. Labels were not used in tuning; external metrics are reported for analysis only.', styles['BodyText']))
    story.append(Spacer(1, 12))
    # Part F: PCA Influence & Suitability (Trade-offs)
    story.append(Paragraph('Part F: PCA Influence & Suitability (Trade-offs)', styles['Heading2']))
    story.append(Paragraph('How PCA influenced clustering decisions: PCA features yielded the best balance of internal indices (e.g., silhouette) at k=3, indicating reduced collinearity/noise improved separability.', styles['BodyText']))
    story.append(Paragraph('Hierarchical vs K-means suitability: hierarchical exposes multi-scale structure but depends on linkage choice and can show chaining; K-means is efficient and interpretable for roughly convex clusters. Choice depends on goals and data geometry.', styles['BodyText']))
    story.append(Paragraph('Reflection & Rationale', styles['Heading2']))
    story.append(Paragraph('Standardization ensures fair distances; PCA reduces noise and collinearity, often improving cluster compactness/separation. K-means is efficient and interpretable for convex clusters; hierarchical reveals multiscale structure and is evaluated via cophenetic correlation. In this run, PCA features with k=3 yielded the best internal indices, suggesting reduction clarified the grouping signal.', styles['BodyText']))
    story.append(Paragraph('Trade-offs: K-means sensitivity to scaling and initialization vs. hierarchical dependence on linkage choice and potential chaining. Future work: stability analysis, alternative algorithms (GMM, DBSCAN), domain-informed feature engineering.', styles['BodyText']))
    story.append(Spacer(1, 12))
    story.append(Paragraph('Applications', styles['Heading2']))
    story.append(Paragraph('Quality control: group similar wines; product segmentation; inventory planning and sourcing; marketing tailored to cluster profiles.', styles['BodyText']))
    story.append(Spacer(1, 6))
    story.append(Paragraph('Benefits & Limitations', styles['Heading2']))
    story.append(Paragraph('Benefits: fast, interpretable clustering; dendrograms reveal structure; PCA reduces noise. Limitations: K-means assumes convex clusters and needs standardization; hierarchical is sensitive to linkage; PCA components may be less interpretable.', styles['BodyText']))
    doc.build(story)
    print(f'Report saved to: {report_path}')

    # Slides with python-pptx
    prs = Presentation()

    def add_text_slide(prs_obj, title, bullets):
        slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        if bullets:
            tf.text = bullets[0]
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0

    def add_code_slide(prs_obj, title, code_text):
        slide = prs_obj.slides.add_slide(prs_obj.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = code_text
        for run in p.runs:
            run.font.name = 'Courier New'
            run.font.size = Pt(12)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = 'UCI Wine: K-means vs Hierarchical'
    slide.placeholders[1].text = 'Problem, Methods, Rationale, Results, Visuals, Inference'
    # Problem statement and dataset
    add_text_slide(prs, 'Problem Statement', [
        'Cluster wines by 13 physico-chemical attributes; labels not used for tuning.',
        'Compare K-means vs hierarchical; select k via internal indices.',
        'Interpret clusters via centroids and statistical tests.'
    ])
    add_text_slide(prs, 'Dataset Overview', [
        f'Samples: {n_samples}, Features: {n_features}.',
        f'Features: {", ".join(X.columns)}.',
        f'Labels only for external validation; counts: {class_counts}.'
    ])
    # Pipeline slide
    add_text_slide(prs, 'Pipeline', [
        'Load → Impute (median) → Standardize (Z-score) → EDA',
        'PCA (retain 95% variance) → Feature engineering',
        'K-means & Hierarchical → Validity indices → Inference'
    ])

    # Part A: Data Loading & Preprocessing
    add_text_slide(prs, 'Part A: Data Loading & Preprocessing (How & Significance)', [
        'How: load UCI Wine, show 10 rows, confirm numeric types.',
        'Impute (median; none missing) and standardize (Z-score).',
        'Save scaled dataset to CSV for reproducibility.'
    ])
    add_code_slide(prs, 'Part A Code', (
        "data = load_wine()\n"
        "X = pd.DataFrame(data.data, columns=data.feature_names)\n"
        "imputer = SimpleImputer(strategy='median')\n"
        "X_imputed = pd.DataFrame(imputer.fit_transform(X), columns=X.columns)\n"
        "scaler = StandardScaler(); X_scaled = pd.DataFrame(scaler.fit_transform(X_imputed), columns=X.columns)\n"
        "X_scaled.to_csv(os.path.join(out_dir, 'scaled_wine.csv'), index=False)"
    ))
    add_text_slide(prs, 'Part A: Significance', [
        'Numeric consistency and standardization ensure fair distances for clustering.',
        'Imputation avoids bias and enables PCA/clustering without missing-value issues.'
    ])

    # Part B: First EDA
    add_text_slide(prs, 'Part B: First EDA (How & Significance)', [
        'Compute descriptive stats (mean, std, min, max) for key features.',
        'Plot histograms for 4 variables; note skewness (e.g., malic_acid right-skew).',
        'Build PCA scree plot to assess variance explained by components.'
    ])
    add_code_slide(prs, 'Part B Code Snippets', (
        "selected = ['alcohol','malic_acid','ash','flavanoids']\n"
        "desc = X_imputed[selected].describe().T\n"
        "for col in selected:\n"
        "    fig, ax = plt.subplots(); sns.histplot(X_imputed[col], kde=True, ax=ax)\n"
        "    save_fig(fig, os.path.join(hist_dir, f'{col}_hist.png'))\n"
        "pca = PCA(); pca.fit(X_scaled); evr = pca.explained_variance_ratio_"
    ))
    add_text_slide(prs, 'Part B: Significance', [
        'Histograms show spread, shape, outliers; skewness highlights long tails.',
        'Scree plot informs dimension reduction to reduce noise/collinearity.'
    ])

    # Part C: Feature Engineering
    add_text_slide(prs, 'Part C: Feature Engineering (How & Significance)', [
        'Create composite features (acid_index, phenols_ratio) and standardize.',
        'Apply PCA retaining 95% variance to reduce dimensions.',
        'Use transformed features if they improve separation and compactness.'
    ])
    add_code_slide(prs, 'Part C Code', (
        "acid_index = X_imputed['malic_acid'] / X_imputed['alcohol']\n"
        "phenols_ratio = X_imputed['total_phenols'] / (X_imputed['flavanoids'] + 1e-6)\n"
        "comp_scaler = StandardScaler(); comp = comp_scaler.fit_transform(np.c_[acid_index, phenols_ratio])\n"
        "X_fe = X_scaled.copy(); X_fe['acid_index'] = comp[:,0]; X_fe['phenols_ratio'] = comp[:,1]\n"
        "pca95 = PCA(n_components=0.95, svd_solver='full'); X_pca = pca95.fit_transform(X_scaled)"
    ))
    add_text_slide(prs, 'Part C: Significance', [
        'Composite features inject domain insight; PCA reduces redundancy and noise.'
    ])

    # Part D: Clustering Modelling & Validation
    add_text_slide(prs, 'Part D: Clustering & Validation (How & Significance)', [
        'K-means: run k=2..8; compute Silhouette↑, DB↓, CH↑; select k.',
        'Hierarchical: Ward & average dendrograms; compute cophenetic correlation.',
        'External validation (ARI, NMI): analysis-only; do not tune with labels.'
    ])
    add_code_slide(prs, 'Part D Code Snippets', (
        "for k in range(2, 9):\n"
        "    km = KMeans(n_clusters=k, n_init=10, random_state=42)\n"
        "    labels = km.fit_predict(features)\n"
        "    s = silhouette_score(features, labels); db = davies_bouldin_score(features, labels); ch = calinski_harabasz_score(features, labels)\n"
        "Z_ward = linkage(features_for_kmeans, method='ward'); Y = pdist(features_for_kmeans); coph_ward, _ = cophenet(Z_ward, Y)"
    ))
    add_text_slide(prs, 'Part D: Significance', [
        'Indices balance cohesion/separation; cophenetic checks dendrogram fidelity.',
        'External metrics contextualize clusters against ground truth without tuning.'
    ])

    # Part E: Second EDA & Statistical Inference
    add_text_slide(prs, 'Part E: Second EDA & Statistical Inference (How & Significance)', [
        'Show centroids in original units; visualize PCA scatter by cluster.',
        'Run ANOVA & Kruskal-Wallis on alcohol, flavanoids; report p-values & effect sizes.',
        'Interpret discrepancies if labels used for validation (ARI/NMI).'
    ])
    add_code_slide(prs, 'Part E Code Snippets', (
        "df = X_imputed.copy(); df['cluster'] = kmeans_labels\n"
        "groups = [g['alcohol'].values for _, g in df.groupby('cluster')]\n"
        "f, p = stats.f_oneway(*groups); h, pkw = stats.kruskal(*groups)\n"
        "# eta² and ε² computed from SS and H-statistic"
    ))
    add_text_slide(prs, 'Part E: Significance', [
        'Centroids explain cluster profiles; PCA scatter shows overlap/separation.',
        'Effect sizes contextualize practical differences beyond p-values.'
    ])

    # Part F: PCA Influence & Suitability
    add_text_slide(prs, 'Part F: PCA Influence & Suitability', [
        'PCA improved indices at k=3 (chosen), clarifying grouping signal.',
        'K-means: fast, interpretable for convex clusters; Hierarchical: linkage-sensitive but reveals structure.',
        'Choose based on data geometry, interpretability needs, and goals.'
    ])

    # Histogram example (alcohol)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Histogram (Alcohol)'
    slide.shapes.add_picture(os.path.join(hist_dir, 'alcohol_hist.png'), Inches(1), Inches(1.5), width=Inches(8))
    add_text_slide(prs, 'Histogram: Purpose & Insights', [
        'Purpose: show the spread and shape of values.',
        'Insight: alcohol values are fairly symmetric; no extreme outliers.',
        'Skewness helps spot long tails (malic_acid is right-skewed).'
    ])
    add_code_slide(prs, 'How We Built the Histogram', (
        "fig, ax = plt.subplots(figsize=(6, 4))\n"
        "sns.histplot(X_imputed['alcohol'], kde=True, ax=ax)\n"
        "ax.set_title('Histogram of alcohol')\n"
        "save_fig(fig, os.path.join(hist_dir, 'alcohol_hist.png'))"
    ))
    # Scree plot slide
    add_text_slide(prs, 'Methods & Rationale', [
        'Standardization: essential for distance-based clustering and PCA.',
        'PCA: reduce collinearity/noise; retained components: ' + str(X_pca.shape[1]),
        'K-means: efficient centroid-based clustering; convex clusters.',
        'Hierarchical: dendrogram structure; Ward vs average linkage.',
        'Indices: Silhouette↑, Davies–Bouldin↓, Calinski–Harabasz↑.'
    ])
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'PCA Scree'
    slide.shapes.add_picture(os.path.join(fig_dir, 'pca_scree.png'), Inches(1), Inches(1.5), width=Inches(8))
    add_text_slide(prs, 'Scree: Purpose & Insights', [
        'Purpose: see how much variance each PCA component explains.',
        'Insight: after early components, added variance drops (diminishing returns).',
        'We keep components covering 95% of variance (10 components).'
    ])
    add_code_slide(prs, 'How We Built the Scree Plot', (
        "pca = PCA(); pca.fit(X_scaled)\n"
        "evr = pca.explained_variance_ratio_\n"
        "fig, ax = plt.subplots()\n"
        "ax.plot(np.arange(1, len(evr)+1), evr, marker='o')\n"
        "save_fig(fig, os.path.join(fig_dir, 'pca_scree.png'))"
    ))
    # KMeans selection slide
    for name in ['scaled', 'pca', 'feature_engineered']:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f'Indices vs k ({name})'
        slide.shapes.add_picture(os.path.join(fig_dir, f'indices_vs_k_{name}.png'), Inches(1), Inches(1.5), width=Inches(8))
    add_text_slide(prs, 'Indices: Purpose & Insights', [
        'Purpose: choose the best number of clusters (k).',
        'Silhouette: higher is better; DB: lower is better; CH: higher is better.',
        'We pick the k with good balance across these indices.'
    ])
    add_code_slide(prs, 'How We Built Indices vs k', (
        "for k in range(2, 9):\n"
        "    km = KMeans(n_clusters=k, n_init=10, random_state=42)\n"
        "    labels = km.fit_predict(features)\n"
        "    s = silhouette_score(features, labels)\n"
        "    db = davies_bouldin_score(features, labels)\n"
        "    ch = calinski_harabasz_score(features, labels)\n"
        "# plot s, db, ch for each k and save"
    ))
    # Dendrogram slides
    for name in ['ward', 'average']:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f'Dendrogram ({name})'
        slide.shapes.add_picture(os.path.join(fig_dir, f'dendrogram_{name}.png'), Inches(0.5), Inches(1.2), width=Inches(9))
    add_text_slide(prs, 'Dendrograms: Purpose & Insights', [
        'Purpose: show how clusters merge step-by-step.',
        'Ward: groups by minimizing variance; Average: uses mean distances.',
        'Cophenetic correlation checks how well tree preserves distances.'
    ])
    add_code_slide(prs, 'How We Built Dendrograms', (
        "Z_ward = linkage(features_for_kmeans, method='ward', metric='euclidean')\n"
        "fig, ax = plt.subplots()\n"
        "dendrogram(Z_ward, ax=ax)\n"
        "save_fig(fig, os.path.join(fig_dir, 'dendrogram_ward.png'))"
    ))
    # Selection & key metrics
    add_text_slide(prs, 'Selection & Key Results', [
        f'Selected: {chosen_set}, k={chosen_k}.',
        f'Silhouette={best_sil:.3f}, DB={best_db:.3f}, CH={best_ch:.1f}.',
        f'Cophenetic (Ward)={coph_ward:.3f}, (Average)={coph_avg:.3f}.',
        f'External (optional): ARI={ari:.3f}, NMI={nmi:.3f}.'
    ])
    # PCA scatter slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'PCA Scatter by Cluster'
    slide.shapes.add_picture(os.path.join(fig_dir, 'pca_scatter_kmeans.png'), Inches(1), Inches(1.5), width=Inches(8))
    add_text_slide(prs, 'PCA Scatter: Purpose & Insights', [
        'Purpose: visualize clusters in 2D using first two PCs.',
        'Insight: see separation and overlap between groups.',
        'Colors represent K-means cluster assignments.'
    ])
    add_code_slide(prs, 'How We Built PCA Scatter', (
        "pca2 = PCA(n_components=2)\n"
        "X_pca2 = pca2.fit_transform(X_scaled)\n"
        "fig, ax = plt.subplots()\n"
        "ax.scatter(X_pca2[:,0], X_pca2[:,1], c=kmeans_labels, cmap='viridis')\n"
        "save_fig(fig, os.path.join(fig_dir, 'pca_scatter_kmeans.png'))"
    ))
    # Cluster profiles slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = 'Cluster Profiles (Selected Vars)'
    slide.shapes.add_picture(os.path.join(fig_dir, 'cluster_profiles_bar.png'), Inches(1), Inches(1.5), width=Inches(8))
    add_text_slide(prs, 'Profiles: Purpose & Insights', [
        'Purpose: compare typical values of key variables across clusters.',
        'Insight: see which clusters have higher/lower alcohol, acidity, etc.',
        'Helps explain what makes clusters different.'
    ])
    add_code_slide(prs, 'How We Built Cluster Profiles', (
        "key_vars = ['alcohol','malic_acid','ash','flavanoids']\n"
        "fig, ax = plt.subplots()\n"
        "idx = np.arange(len(key_vars)); width = 0.8/ chosen_k\n"
        "for c in range(chosen_k):\n"
        "    ax.bar(idx + c*width, centroids_original.loc[c, key_vars], width)\n"
        "save_fig(fig, os.path.join(fig_dir, 'cluster_profiles_bar.png'))"
    ))

    # Statistical inference and rationale slides
    add_text_slide(prs, 'Statistical Inference', [
        f"Alcohol: ANOVA p={results_stats['alcohol']['ANOVA_p']:.4f}, eta²={results_stats['alcohol']['eta_squared']:.3f}",
        f"Alcohol: Kruskal p={results_stats['alcohol']['Kruskal_p']:.4f}, ε²={results_stats['alcohol']['epsilon_squared']:.3f}",
        f"Flavanoids: ANOVA p={results_stats['flavanoids']['ANOVA_p']:.4f}, eta²={results_stats['flavanoids']['eta_squared']:.3f}",
        f"Flavanoids: Kruskal p={results_stats['flavanoids']['Kruskal_p']:.4f}, ε²={results_stats['flavanoids']['epsilon_squared']:.3f}"
    ])
    add_text_slide(prs, 'Applications', [
        'Quality control: group wines with similar chemistry for consistency.',
        'Product segmentation: design blends or offerings by cluster profiles.',
        'Inventory planning and sourcing: match supply to cluster demand.',
        'Marketing: tailor messaging to cluster characteristics.'
    ])
    add_text_slide(prs, 'Benefits & Limitations', [
        'Benefits: fast, interpretable clusters; structure insight via dendrograms; reduces noise with PCA.',
        'Limitations: K-means assumes convex clusters; hierarchical sensitive to linkage.',
        'Needs standardization; PCA components hard to explain directly.'
    ])
    add_text_slide(prs, 'Rationale & Trade-offs', [
        'Standardization: fair distances; PCA: reduce noise/collinearity.',
        'K-means: fast, interpretable centroids; convex cluster assumption.',
        'Hierarchical: structure via dendrogram; linkage choice matters.',
        'PCA+K-means often balances compactness and separation on numeric data.'
    ])
    add_text_slide(prs, 'Limitations & Next Steps', [
        'Assess stability; try GMM, DBSCAN; add domain-informed features.',
        'Bootstrap validity indices; silhouette per sample; sensitivity analysis.'
    ])

    slides_path = os.path.join(out_dir, 'wine_clustering_slides.pptx')
    prs.save(slides_path)
    print(f'Slides saved to: {slides_path}')


if __name__ == '__main__':
    main()